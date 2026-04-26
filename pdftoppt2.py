
import streamlit as st
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import os

st.set_page_config(
    page_title="Free PDF to PPT Converter | No Upload Limits | Secure & Local",
    page_icon="📄",
    menu_items={
        'Get Help': 'https://your-contact-page.com',
        'About': "# This tool converts PDF to editable PowerPoint locally in your browser."
    }
)

# --- UI Setup ---
st.set_page_config(page_title="Journalist PDF Converter", page_icon="📄")
st.title("🔒 Secure PDF to PowerPoint")
st.markdown("This tool processes files **locally** on your MacBook M1.")

uploaded_file = st.file_uploader("Select the secret PDF file", type="pdf")

if uploaded_file:
    # Save uploaded file locally
    with open(uploaded_file.name, "wb") as f:
        f.write(uploaded_file.getbuffer())

    st.success(f"File '{uploaded_file.name}' ready.")

    if st.button("Convert to PPTX"):
        output_name = uploaded_file.name.replace(".pdf", ".pptx")

        with st.spinner("Converting... This may take a moment for large files."):
            try:
                # 1. Convert PDF to Images (High Res for readability)
                # This uses the Poppler you installed via Homebrew
                images = convert_from_path(uploaded_file.name, dpi=200)

                # 2. Create PowerPoint
                prs = Presentation()

                # Use Widescreen dimensions
                prs.slide_width = Inches(13.33)
                prs.slide_height = Inches(7.5)

                for i, image in enumerate(images):
                    temp_img_path = f"temp_page_{i}.png"
                    image.save(temp_img_path, "PNG")

                    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
                    slide.shapes.add_picture(temp_img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

                    os.remove(temp_img_path)  # Clean up as we go

                prs.save(output_name)

                # 3. Download Button
                with open(output_name, "rb") as file:
                    st.download_button(
                        label="📥 Download converted PPTX",
                        data=file,
                        file_name=output_name,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                st.balloons()

                # Cleanup original upload
                os.remove(uploaded_file.name)

            except Exception as e:
                st.error(f"Processing Error: {e}")
                st.info("Tip: Ensure 'brew install poppler' was successful in your Terminal.")
